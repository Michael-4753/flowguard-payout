#!/usr/bin/env python3
# Generates FlowGuard-Hackathon-Pitch.pptx (16:9, Chinese) into public/.
# Uses python-pptx; embeds the demo illustration; CJK-safe via East-Asian font.
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

ROOT = "/home/user/flowguard-payout"
OUT = os.path.join(ROOT, "public", "FlowGuard-Hackathon-Pitch.pptx")
DEMO = os.path.join(ROOT, "public", "pitch-shots", "demo.png")
CN = "Noto Sans CJK SC"

TEAL = RGBColor(0x0F, 0x76, 0x6E)
INK = RGBColor(0x0F, 0x17, 0x2A)
MUT = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARDBG = RGBColor(0xF1, 0xF5, 0xF9)
RED = RGBColor(0x99, 0x1B, 0x1B)
REDBG = RGBColor(0xB9, 0x1B, 0x1B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def set_cn(run, bold=False):
    run.font.name = CN
    r = run._r
    rpr = r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        e = rpr.find(qn(tag))
        if e is None:
            e = rpr.makeelement(qn(tag), {}); rpr.append(e)
        e.set("typeface", CN)
    run.font.bold = bold

def add_text(slide, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.color.rgb = color
    set_cn(run, bold)
    return tb

def rect(slide, x, y, w, h, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp

def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def eyebrow(slide, text, color=TEAL):
    add_text(slide, Inches(0.9), Inches(0.55), Inches(11.5), Inches(0.4), text, 14, color, bold=True)

def title(slide, text, color=INK, size=36):
    add_text(slide, Inches(0.9), Inches(0.95), Inches(11.5), Inches(1.0), text, size, color, bold=True)

def subtitle(slide, text, color=MUT):
    add_text(slide, Inches(0.9), Inches(1.95), Inches(11.5), Inches(0.8), text, 16, color)

def bullets_grid(slide, items, top=2.7, card_fill=CARDBG, head_color=INK, body_color=MUT, cols=2, ch=1.25):
    gap = 0.25
    total_w = 11.5
    cw = (total_w - gap * (cols - 1)) / cols
    for idx, (head, body) in enumerate(items):
        r = idx // cols; c = idx % cols
        x = Inches(0.9 + c * (cw + gap)); y = Inches(top + r * (ch + gap))
        rect(slide, x, y, Inches(cw), Inches(ch), card_fill)
        add_text(slide, x + Inches(0.2), y + Inches(0.12), Inches(cw - 0.4), Inches(0.4), head, 16, head_color, bold=True)
        add_text(slide, x + Inches(0.2), y + Inches(0.55), Inches(cw - 0.4), Inches(ch - 0.65), body, 12.5, body_color)

def footnote(slide, text, color=MUT):
    add_text(slide, Inches(0.9), Inches(6.9), Inches(11.5), Inches(0.4), text, 11, color)

# ---------- 1 cover ----------
s = prs.slides.add_slide(BLANK); bg(s, TEAL)
add_text(s, Inches(0.9), Inches(1.6), Inches(11.5), Inches(0.5),
         "2026 创青春 AI 黑客松 · 顺德行 · 自由创新赛道", 15, RGBColor(0xCC,0xEE,0xEA), bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(1.4), "FlowGuard", 72, WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.4), Inches(3.9), Inches(10.5), Inches(1.0),
         "把 AI 带进真实场景：跨境付款「先核查，再付款」的风控控制台", 22, RGBColor(0xEA,0xF7,0xF5), align=PP_ALIGN.CENTER)
add_text(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.5),
         "一人团队 · 负责 App 全部内容 · 5 分钟路演", 13, RGBColor(0xBF,0xE4,0xE0), align=PP_ALIGN.CENTER)

# ---------- 2 toc ----------
s = prs.slides.add_slide(BLANK); bg(s, WHITE)
eyebrow(s, "目录"); title(s, "今天用 5 分钟讲清这几件事")
toc = ["真实案例：一笔货款怎么踩坑","真实问题：谁在痛、痛在哪","产品介绍：FlowGuard 是什么",
       "三大核心功能：预检 · 链路透明 · 对账看板",
       "AI 覆盖一览：四痛点→四能力","Demo 视频：逐页看页面功能","技术架构：怎么实现的","合规边界：我们不碰什么（关键）","总结与可继续方向"]
for idx, item in enumerate(toc):
    c = idx % 3; r = idx // 3
    x = Inches(0.9 + c*3.95); y = Inches(2.55 + r*1.35)
    rect(s, x, y, Inches(3.75), Inches(1.15), CARDBG)
    add_text(s, x+Inches(0.15), y+Inches(0.12), Inches(0.5), Inches(0.5), str(idx+1), 18, TEAL, bold=True)
    add_text(s, x+Inches(0.7), y+Inches(0.2), Inches(2.9), Inches(0.8), item, 13, INK, bold=True)

# ---------- 3 case ----------
s = prs.slides.add_slide(BLANK); bg(s, WHITE)
eyebrow(s, "先看一笔真实的付款")
title(s, "一笔 $48,000 越南货款，是怎么踩进坑里的", size=30)
subtitle(s, "周一，出纳给越南供应商汇一笔货款。看似普通的一笔电汇，一周内连踩四个坑——每个坑，正是下一页要讲的痛点。")

# 左：时间线叙事
rect(s, Inches(0.9), Inches(2.75), Inches(5.55), Inches(3.95), CARDBG)
add_text(s, Inches(1.15), Inches(2.9), Inches(5.1), Inches(0.4), "一周时间线", 15, TEAL, bold=True)
timeline = [
    ("周一", "出纳发起 $48,000 电汇，钱汇出，开始等待。"),
    ("第 3 天", "供应商说没收到。问银行，只答「在途」，卡在哪家中间行没人说得清。"),
    ("第 6 天", "银行退回：收款公司名与账户不符（一个字母之差）。$133 手续费打水漂、货期延误。"),
    ("第 7 天", "想改走别的通道，却不清楚越南走廊要哪些材料；对账时退回款/手续费/二次汇款三笔对不上。"),
]
ty = 3.35
for when, what in timeline:
    add_text(s, Inches(1.15), Inches(ty), Inches(1.1), Inches(0.4), when, 12.5, INK, bold=True)
    add_text(s, Inches(2.3), Inches(ty), Inches(3.95), Inches(0.8), what, 11.5, MUT)
    ty += 0.83

# 右：四个踩坑点 → 预告四痛点
rect(s, Inches(6.7), Inches(2.75), Inches(5.7), Inches(3.95), RGBColor(0xFB,0xEC,0xEC))
add_text(s, Inches(6.95), Inches(2.9), Inches(5.2), Inches(0.4), "四个坑，四个痛点", 15, RGBColor(0xC0,0x3A,0x2B), bold=True)
pits = [
    ("看不到钱卡在哪", "全程黑箱，只能干等 → 痛点①"),
    ("发出后才知会被退", "收款信息问题事后才暴露 → 痛点②"),
    ("走廊规则要重新摸", "每国每行要求都不同 → 痛点③"),
    ("三笔款对不上账", "退回/手续费/重汇难核销 → 痛点④"),
]
py = 3.35
for head, body in pits:
    add_text(s, Inches(6.95), Inches(py), Inches(5.2), Inches(0.35), head, 13, RGBColor(0xC0,0x3A,0x2B), bold=True)
    add_text(s, Inches(6.95), Inches(py+0.35), Inches(5.2), Inches(0.4), body, 11.5, MUT)
    py += 0.83
footnote(s, "这不是虚构极端案例，而是外贸出纳每周都在经历的常态。下一页，把它拆成四个可被验证的痛点。")

# ---------- 4 problem ----------
s = prs.slides.add_slide(BLANK); bg(s, WHITE)
eyebrow(s, "真实问题 · 谁会使用"); title(s, "跨境 B2B 付款，至今仍是「盲发」")
subtitle(s, "使用者：外贸企业出纳/财务、跨境电商、做海外供应商结算的中小企业。下面四痛点，正好对应后面四个 AI 能力。")
bullets_grid(s, [
    ("① 中间行黑箱，钱卡哪看不见","看不到钱走了哪条路、卡在哪个中间行，只能干等 → AI: where is my money?"),
    ("② 莫名被退，发起时才知道","被退回的以转账（尤其打款给个人收款人）最常见；退回多是收款信息/合规校验问题，却发起后才发现。传统核实要银行→出纳→业务→供应商层层传话，拖几天到一周 → AI: draft the follow-up"),
    ("③ 多国供应商，重复适应银行","多国多供应商，每次结算都要重新适应不同银行要求，效率极低 → AI: what this corridor requires"),
    ("④ 财务对账，凭证对不上","链上/链下凭证对不上，多笔跨国付款时尤其头疼 → AI: why don't these match?"),
], ch=1.5)
footnote(s, "真实、高频、可被验证的痛点——每一条都在后续被一个 AI 能力接住。")

# ---------- 4 product ----------
s = prs.slides.add_slide(BLANK); bg(s, WHITE)
eyebrow(s, "产品介绍 · 解决什么"); title(s, "付款前的「退回风险 + 合规」控制台")
subtitle(s, "一句话：在把付款指令提交给持牌机构之前，先预检退回风险、比价选路、双人放行（大额可升级为多级/多签）。")
bullets_grid(s, [
    ("全链路闭环","预检→供应商核实→双人审批（大额可升级多级/多签）→生成付款指令(提交持牌机构)→到账确认→自动对账。"),
    ("多通道比价","统一池管理多类持牌结算通道，按费用/时效/退回率自动排序，给出可解释的最优路径推荐。"),
    ("面向真实用户","为中小外贸出纳设计，移动端优先，几步走完一笔高风险付款的完整控制。"),
    ("产出是指令","批准即生成付款指令，交由持牌机构划付——平台不经手资金。"),
])

# ---------- 5/6/7 features ----------
def feature(eb, tt, st, items, shot=None, video=None):
    s = prs.slides.add_slide(BLANK); bg(s, WHITE)
    eyebrow(s, eb); title(s, tt, size=30)
    add_text(s, Inches(0.9), Inches(1.7), Inches(7.7), Inches(0.7), st, 14, MUT)
    lw = 7.7; ch = 1.02; gap = 0.18; top = 2.6
    for idx, (head, body) in enumerate(items):
        y = Inches(top + idx * (ch + gap))
        rect(s, Inches(0.9), y, Inches(lw), Inches(ch), CARDBG)
        add_text(s, Inches(1.1), y + Inches(0.12), Inches(lw - 0.4), Inches(0.4), f"{idx+1}. {head}", 15, INK, bold=True)
        add_text(s, Inches(1.1), y + Inches(0.5), Inches(lw - 0.4), Inches(0.5), body, 11.5, MUT)
    # Right column: embed a pre-recorded demo video when provided; otherwise keep
    # an empty placeholder box (same region/size) so the layout stays intact and
    # a screenshot/video can be added manually later. `shot` kept for compat.
    if video and os.path.exists(video):
        rect(s, Inches(8.9), Inches(1.65), Inches(3.55), Inches(5.4), RGBColor(0xF1, 0xF5, 0xF9))
        s.shapes.add_movie(video, Inches(8.98), Inches(1.75), Inches(3.4), Inches(5.2),
                           mime_type="video/mp4")
        add_text(s, Inches(8.9), Inches(7.12), Inches(3.55), Inches(0.3), "Demo 视频（点击播放）", 9, MUT, align=PP_ALIGN.CENTER)
    else:
        rect(s, Inches(8.9), Inches(1.65), Inches(3.55), Inches(5.4), RGBColor(0xF1, 0xF5, 0xF9))
        add_text(s, Inches(8.9), Inches(4.15), Inches(3.55), Inches(0.4), "在此处粘贴截图", 11, MUT, align=PP_ALIGN.CENTER)
        add_text(s, Inches(8.9), Inches(7.12), Inches(3.55), Inches(0.3), "产品界面（现场可实测）", 9, MUT, align=PP_ALIGN.CENTER)
    return s
SHOTS = "/home/user/flowguard-payout/public/pitch-shots"
FEAT1_VIDEO = "/home/user/flowguard-payout/public/eazo-assets/att_70duam7r6rrsf3w7-896ee20d49-feat-precheck-demo-v2.mp4"
FEAT2_VIDEO = "/home/user/flowguard-payout/public/eazo-assets/att_385lhafcz161fcq3-190894a7b5-feat-route-demo.mp4"
FEAT3_VIDEO = "/home/user/flowguard-payout/public/eazo-assets/att_3qmrjx5a4911h6it-3fb6d329d4-feat-reconcile-demo.mp4"
feature("核心功能① · AI 为何必要","收款方信息预检-AI Agent","收款账户信息智能校验、国别规则适配、退回风险评分、合规风险提示报告。",
        [("账户信息智能校验","收款名/SWIFT/IBAN 一致性与账户状态校验，先拦截而非先失败。"),
         ("国别规则适配","按收款人国家/币种自动适配该走廊的结算要求，省去每次重新摸规则。"),
         ("退回风险评分","确定性规则引擎给出退回概率与命中因子；AI 补充语义/情境风险，只加警示、不降评分。"),
         ("合规风险提示报告","把检查项翻译成通俗解释与修复步骤，并可一键生成供应商核实工单。")],
        shot=f"{SHOTS}/feat-precheck.png", video=FEAT1_VIDEO)
feature("核心功能② · 事中破黑箱","结算链路透明化追踪","逐跳看清资金流经的每一层中间行，破除“钱卡在哪看不见”的黑箱；并在多持牌通道间比价选路。",
        [("资金流逐跳可视","从起点到收款人，逐一层展开经手行、耗时与逐层扣费后余额，钱走到哪一目了然。"),
         ("黑箱卡点标注","逐跳标出透明度（清晰/部分/黑箱）与滞留卡点，付款后还能定位钱现在卡在哪家中间行。"),
         ("多通道比价选路","看清链路后在多类持牌结算通道间，按费用/时效/退回率自动比价，给出可解释的最优路径——平台不执行付款。"),
         ("为什么要 AI","AI 用自然语言说清钱现在卡在哪家中间行、为何被卡、还要多久、你能做什么。")],
        shot=f"{SHOTS}/feat-route.png", video=FEAT2_VIDEO)
feature("核心功能③ · 成果可演示","统一结算状态与对账看板","聚合各通道汇款进度、中转链路信息、交易凭证，实现多笔跨国付款台账可视化与自动对账核销。",
        [("多通道进度聚合","汇总各持牌通道的汇款进度与中转链路信息，钱走到哪一目了然。"),
         ("凭证与台账可视化","聚合交易凭证，多笔跨国付款台账集中可视化管理。"),
         ("自动对账核销","应收 vs 实收、费用与汇兑损失自动匹配核销，并可导出对账单。"),
         ("结算链路透明","结算全程可追踪；平台只做透明化追踪，不经手资金。")],
        shot=f"{SHOTS}/feat-reconcile.png", video=FEAT3_VIDEO)

# ---------- 8 AI coverage table ----------
s = prs.slides.add_slide(BLANK); bg(s, WHITE)
eyebrow(s, "AI 为何必要 · 四处真实痛点 → 四个 AI 能力")
title(s, "AI 覆盖一览：不是装饰，是每个卡点的解法", size=30)
add_text(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(0.5),
         "每个 AI 能力都绑定一个真实痛点、一个具体入口、一件明确的事。", 14, MUT)

ai_rows = [
    ("① 看不到钱卡哪", "历史里进行中付款的资金流链路图下方",
     "「AI: where is my money?」", "用自然语言说清钱现在卡在哪家中间行、为何被卡、大概还要多久、你能做什么。"),
    ("② 退回事后才知道", "核验 Case（open 状态）",
     "「AI: draft the follow-up」", "基于命中因子起草给供应商的核实邮件/话术，出纳一键直触供应商、砍掉银行→出纳→业务→供应商的中间层层传话（补齐闭环）。"),
    ("③ 多国重复适应银行", "添加收款人表单（填了国家+银行后出现）",
     "「AI: what this corridor requires」", "生成该国家/币种的结算要求清单（收款名匹配、SWIFT/IBAN、FX 管制文件等），省去每次重新摸规则。"),
    ("④ 对账对不上", "对账卡片（有差异或凭证未匹配时）",
     "「AI: why don't these match?」", "指出金额/凭证对不上的最可能原因 + 下一步。"),
]
tbl_shape = s.shapes.add_table(len(ai_rows) + 1, 3, Inches(0.9), Inches(2.35), Inches(11.5), Inches(4.0))
table = tbl_shape.table
table.columns[0].width = Inches(2.1)
table.columns[1].width = Inches(3.2)
table.columns[2].width = Inches(6.2)
# disable banded default style header emphasis; set our own colors
hdrs = ["痛点", "入口", "AI 做什么"]
for c, htext in enumerate(hdrs):
    cell = table.cell(0, c)
    cell.fill.solid(); cell.fill.fore_color.rgb = TEAL
    cell.margin_left = Inches(0.1); cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)
    p = cell.text_frame.paragraphs[0]
    run = p.add_run(); run.text = htext
    run.font.size = Pt(13); run.font.color.rgb = WHITE; set_cn(run, True)
for r, (pain, entry, call, does) in enumerate(ai_rows, start=1):
    row_fill = RGBColor(0xF8, 0xFA, 0xFC) if r % 2 else WHITE
    cells_text = [(pain, INK, True, 12), (entry, MUT, False, 11.5), (None, None, None, None)]
    for c, val in enumerate(cells_text[:2]):
        text, color, bold, sz = val
        cell = table.cell(r, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = row_fill
        cell.margin_left = Inches(0.1); cell.margin_top = Inches(0.05); cell.margin_bottom = Inches(0.05)
        p = cell.text_frame.paragraphs[0]
        run = p.add_run(); run.text = text
        run.font.size = Pt(sz); run.font.color.rgb = color; set_cn(run, bold)
    # third cell: call (teal bold) + does (muted)
    cell = table.cell(r, 2)
    cell.fill.solid(); cell.fill.fore_color.rgb = row_fill
    cell.margin_left = Inches(0.1); cell.margin_top = Inches(0.05); cell.margin_bottom = Inches(0.05)
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = call; r1.font.size = Pt(12); r1.font.color.rgb = TEAL; set_cn(r1, True)
    r2 = p.add_run(); r2.text = " — " + does; r2.font.size = Pt(11.5); r2.font.color.rgb = MUT; set_cn(r2, False)
add_text(s, Inches(0.9), Inches(6.95), Inches(11.5), Inches(0.4),
         "AI 只加价值、只加警示——绝不降低确定性引擎评分；输出均需提交指令前人工核实。", 10.5, MUT)

# ---------- 9 demo ----------
s = prs.slides.add_slide(BLANK); bg(s, WHITE)
eyebrow(s, "预录 Demo 视频 · 逐页展示"); title(s, "Demo 视频：逐页看每个页面在做什么", size=32)
add_text(s, Inches(0.9), Inches(1.85), Inches(11.5), Inches(0.6),
         "依次走过每个页面：首页概览 → 退回风险预检 → 多通道比价选路 → 结算对账看板 → 双人放行审批 → AI 核实工单。", 15, MUT)
# Screenshot removed — keep an empty placeholder box (same region) so the slide
# layout stays intact and a demo video can be placed in manually later.
rect(s, Inches(1.35), Inches(2.5), Inches(10.6), Inches(4.3), RGBColor(0xF1, 0xF5, 0xF9))
add_text(s, Inches(1.35), Inches(4.5), Inches(10.6), Inches(0.5), "在此处放入 Demo 视频", 14, MUT, align=PP_ALIGN.CENTER)
footnote(s, "预录 Demo 视频、逐页展示各功能页面——不是 PPT 概念。可扫码访问在线应用实测。")

# ---------- 9 architecture ----------
s = prs.slides.add_slide(BLANK); bg(s, WHITE)
eyebrow(s, "过程可说明 · 技术架构"); title(s, "怎么实现的")
subtitle(s, "Next.js(App Router)+TypeScript+Tailwind；确定性规则引擎+DeepSeek AI；双语 i18n。")
bullets_grid(s, [
    ("前端 / 应用","Next.js 15 + React + Tailwind，移动端优先；react-i18next 中英双语。"),
    ("风控内核","确定性退回风险规则引擎(可解释、可复现) + DeepSeek 作为补充语义信号层。"),
    ("数据与流转","付款指令/收款人台账/工单/到账确认/对账 全链路数据模型。"),
    ("工程质量","全站 TypeScript 0 error、双语键 415/415 对齐、关键路由端到端可跑。"),
])

# ---------- 10 compliance ----------
s = prs.slides.add_slide(BLANK); bg(s, RED)
add_text(s, Inches(0.9), Inches(0.55), Inches(11.5), Inches(0.4), "合规边界 · 打消最大顾虑", 14, RGBColor(0xFE,0xCA,0xCA), bold=True)
add_text(s, Inches(0.9), Inches(0.95), Inches(11.5), Inches(1.0), "我们「不碰」什么", 40, WHITE, bold=True)
add_text(s, Inches(0.9), Inches(1.95), Inches(11.5), Inches(0.7),
         "FlowGuard 是纯软件决策支持工具——这是产品的底线，也是可持续的前提。", 16, RGBColor(0xFE,0xE2,0xE2))
comp = [("✕ 不持牌","不持有任何支付/金融牌照，不以金融机构身份展业。"),
        ("✕ 不经手资金","不收款、不放款、不托管资金；任何环节都不流经资金。"),
        ("✕ 不做加密兑换/转账","不提供稳定币/加密货币的兑换、托管或转账服务。"),
        ("✕ 由持牌机构结算","资金结算均由持牌金融机构完成；平台只做风控、选路、生成指令与追踪。")]
gap=0.25; cw=(11.5-gap)/2; ch=1.3
for idx,(head,body) in enumerate(comp):
    c=idx%2; r=idx//2
    x=Inches(0.9+c*(cw+gap)); y=Inches(2.75+r*(ch+gap))
    rect(s,x,y,Inches(cw),Inches(ch),REDBG)
    add_text(s,x+Inches(0.2),y+Inches(0.12),Inches(cw-0.4),Inches(0.4),head,17,WHITE,bold=True)
    add_text(s,x+Inches(0.2),y+Inches(0.58),Inches(cw-0.4),Inches(0.6),body,13,RGBColor(0xFE,0xE2,0xE2))
add_text(s, Inches(0.9), Inches(6.85), Inches(11.5), Inches(0.4),
         "全站措辞已按此口径统一，并附双语《合规审查报告》PDF。", 11, RGBColor(0xFE,0xCA,0xCA))

# ---------- 11 summary ----------
s = prs.slides.add_slide(BLANK); bg(s, TEAL)
add_text(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(0.4), "总结 · 方向可继续", 14, RGBColor(0xCC,0xEE,0xEA), bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(1.0), "真实问题 · AI 必要 · 现场可演示 · 边界清晰", 34, WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.4), Inches(2.7), Inches(10.5), Inches(0.7),
         "一人团队，独立完成 App 全部内容：产品、风控逻辑、AI 集成、双语、合规措辞与文档。", 16, RGBColor(0xEA,0xF7,0xF5), align=PP_ALIGN.CENTER)
cards=[("已验证","完整闭环可现场跑通；工程质量与合规口径均已落地。"),
       ("可继续","接入更多真实持牌通道、只读 ERP 对账、案例库随使用增长。"),
       ("一句话","让每一笔跨境付款，在提交给持牌机构之前，先被 AI 与规则一起看一眼。")]
cw=3.7; gap=0.3
for idx,(h,b) in enumerate(cards):
    x=Inches(0.9+idx*(cw+gap)); y=Inches(3.9)
    rect(s,x,y,Inches(cw),Inches(1.9),RGBColor(0x11,0x5E,0x59))
    add_text(s,x+Inches(0.2),y+Inches(0.15),Inches(cw-0.4),Inches(0.4),h,17,WHITE,bold=True)
    add_text(s,x+Inches(0.2),y+Inches(0.65),Inches(cw-0.4),Inches(1.1),b,13,RGBColor(0xDF,0xF1,0xEF))
add_text(s, Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.4),
         "FlowGuard · 2026 创青春 AI 黑客松 · 自由创新赛道", 12, RGBColor(0xBF,0xE4,0xE0), align=PP_ALIGN.CENTER)

prs.save(OUT)
print("PPTX written:", OUT, os.path.getsize(OUT), "bytes | slides:", len(prs.slides._sldIdLst))
